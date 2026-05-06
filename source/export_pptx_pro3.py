from e2e_test import build_storyboard
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C = {
    "navy":   RGBColor(0x1E, 0x27, 0x61),
    "mid":    RGBColor(0x2E, 0x3F, 0x8F),
    "dark":   RGBColor(0x0F, 0x16, 0x40),
    "ice":    RGBColor(0xCA, 0xDC, 0xFC),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "accent": RGBColor(0x5B, 0x8C, 0xFF),
    "muted":  RGBColor(0x88, 0x92, 0xB0),
}

TYPE_BG = {
    "introduction": RGBColor(0x1E, 0x27, 0x61),
    "objectives":   RGBColor(0x2E, 0x3F, 0x8F),
    "screen":       RGBColor(0x1A, 0x6B, 0x3C),
    "cyu":          RGBColor(0xB8, 0x5A, 0x00),
    "summary":      RGBColor(0x6B, 0x2F, 0xA0),
}

TYPE_LABELS = {
    "introduction": "INTRODUCTION",
    "objectives":   "OBJECTIVES",
    "screen":       "CONTENT SCREEN",
    "cyu":          "CHECK YOUR UNDERSTANDING",
    "summary":      "SUMMARY",
}

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, x, y, w, h, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return tb

def add_rect(slide, x, y, w, h, color):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_title_slide(pres, topic):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    set_bg(slide, C["dark"])
    add_rect(slide, 0, 5.1, 10, 0.525, C["navy"])
    add_textbox(slide, "eLEARNING COURSE", 0.5, 0.55, 2.4, 0.38, 9, C["white"], bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, topic, 0.5, 1.1, 9, 1.4, 44, C["white"], bold=True)
    add_textbox(slide, "AI-Generated eLearning Storyboard", 0.5, 2.45, 9, 0.7, 22, C["ice"])
    # Removed subtitle line
    add_textbox(slide, f"{topic}  •  eLearning Series", 0.5, 5.1, 9, 0.525, 11, C["ice"])

def add_content_slide(pres, screen, topic):
    stype = screen["type"]
    bg = TYPE_BG.get(stype, C["navy"])
    label = TYPE_LABELS.get(stype, stype.upper())
    title = screen.get("title", label)

    slide = pres.slides.add_slide(pres.slide_layouts[6])
    set_bg(slide, C["dark"])

    # Header band
    add_rect(slide, 0, 0, 10, 1.2, bg)
    add_textbox(slide, label, 0.4, 0, 9, 0.45, 10, C["ice"], bold=True)
    add_textbox(slide, title, 0.4, 0.38, 9, 0.75, 26, C["white"], bold=True)

    # Content card background
    add_rect(slide, 0.4, 1.35, 9.2, 3.9, C["mid"])

    # Build content
    if stype == "cyu":
        content = f"Q: {screen.get('question', '')}\n\n"
        opts = screen.get("options", {})
        correct = screen.get("correct_answer", "")
        for k, v in opts.items():
            tick = "  ✓" if k == correct else ""
            content += f"{k})  {v}{tick}\n"
        content += f"\nCorrect Answer: {correct}"
    else:
        raw = screen.get("content", "")
        lines = [l for l in raw.splitlines() if not l.strip().lower().startswith("title:")]
        content = "\n".join(lines).strip()

    add_textbox(slide, content, 0.6, 1.5, 8.8, 3.6, 14, C["ice"])
    add_textbox(slide, f"{topic}  •  {label}", 0.4, 5.3, 9.2, 0.3, 9, C["muted"], italic=True)

def build_pptx(storyboard, output_file, topic):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, topic)
    for screen in storyboard:
        add_content_slide(prs, screen, topic)

    prs.save(output_file)
    print(f"✅ Saved: {output_file}")

if __name__ == "__main__":
    jobs = [
        ("sample.pdf",  "output_sample1.pptx", "Artificial Intelligence"),
        ("sample2.pdf", "output_sample2.pptx", "Cybersecurity and Data Privacy"),
        ("sample3.pdf", "output_sample3.pptx", "Business and Digital Marketing"),
    ]
    for pdf, out, topic in jobs:
        print(f"\n📄 Processing: {topic}")
        storyboard = build_storyboard(pdf)
        build_pptx(storyboard, out, topic)
    print("\n🎉 All 3 professional PowerPoint files created!")
