
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from e2e_test import build_storyboard

COLORS = {
    "introduction": RGBColor(0x44, 0x72, 0xC4),  # Blue
    "objectives":   RGBColor(0xFF, 0xD9, 0x66),  # Yellow
    "screen":       RGBColor(0x70, 0xAD, 0x47),  # Green
    "cyu":          RGBColor(0xED, 0x7D, 0x31),  # Orange
    "summary":      RGBColor(0x9B, 0x59, 0xB6),  # Purple
}

TEXT_COLORS = {
    "introduction": RGBColor(0xFF, 0xFF, 0xFF),
    "objectives":   RGBColor(0x00, 0x00, 0x00),
    "screen":       RGBColor(0xFF, 0xFF, 0xFF),
    "cyu":          RGBColor(0xFF, 0xFF, 0xFF),
    "summary":      RGBColor(0xFF, 0xFF, 0xFF),
}

def add_slide(prs, screen_type, title, content):
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background colour
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS.get(screen_type, RGBColor(0xFF, 0xFF, 0xFF))

    text_color = TEXT_COLORS.get(screen_type, RGBColor(0x00, 0x00, 0x00))

    # Title box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = text_color

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(18)
    p2.font.color.rgb = text_color

    return slide

def export_pptx(storyboard: list, filename: str = "storyboard.pptx"):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    for screen in storyboard:
        screen_type = screen["type"]
        title = screen["title"]

        if screen_type == "cyu":
            options = screen.get("options", {})
            content = f"Q: {screen.get('question', '')}\n\n"
            for k, v in options.items():
                content += f"{k})  {v}\n"
            content += f"\n✅ Correct Answer: {screen.get('correct_answer', '')}"
        else:
            content = screen.get("content", "")

        add_slide(prs, screen_type, title, content)

    prs.save(filename)
    print(f"✅ PowerPoint saved as {filename}")

if __name__ == "__main__":
    print("=== Building Storyboard ===\n")
    storyboard = build_storyboard("sample.pdf")
    print("\n=== Exporting to PowerPoint ===\n")
    export_pptx(storyboard)
    print("\n✅ Done! Open storyboard.pptx to check all slides.")
