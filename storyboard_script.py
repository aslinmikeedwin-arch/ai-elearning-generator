import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation

def extract_slides(pptx_path):
    prs = Presentation(pptx_path)
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and "AI Learning Series" not in text:
                    texts.append(text)
        slides_data.append({"slide": i, "content": "\n".join(texts)})
    return slides_data

def build_excel(pptx_path, output_path):
    slides = extract_slides(pptx_path)

    screens = [
        {"screen_label": "screen_1",  "title": "introduction",               "slide_index": 1},
        {"screen_label": "screen_2",  "title": "learning objectives",         "slide_index": 2},
        {"screen_label": "screen_3",  "title": "ai fundamentals",             "slide_index": 3},
        {"screen_label": "screen_4",  "title": "check your understanding",    "slide_index": 4},
        {"screen_label": "screen_5",  "title": "what is artificial intelligence?", "slide_index": 5},
        {"screen_label": "screen_6",  "title": "check your understanding",    "slide_index": 6},
        {"screen_label": "screen_7",  "title": "machine learning basics",     "slide_index": 7},
        {"screen_label": "screen_8",  "title": "check your understanding",    "slide_index": 8},
        {"screen_label": "screen_9",  "title": "summary",                     "slide_index": 9},
    ]

    HEADER_BG = "F4CCCC"
    HEADER_FG = "CC0000"

    def fill(color):
        return PatternFill(start_color=color, end_color=color, fill_type="solid")

    def thin_border():
        s = Side(style="thin", color="CCCCCC")
        return Border(left=s, right=s, top=s, bottom=s)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Storyboard"
    ws.sheet_view.showGridLines = False

    ws.column_dimensions["A"].width = 25
    ws.column_dimensions["B"].width = 55
    ws.column_dimensions["C"].width = 45
    ws.column_dimensions["D"].width = 18

    current_row = 1

    for screen in screens:
        audio = slides[screen["slide_index"]]["content"] if screen["slide_index"] < len(slides) else ""
        ost = f"Key points from: {screen['title'].title()}"

        # Screen label row
        ws.row_dimensions[current_row].height = 18
        cell = ws.cell(row=current_row, column=2, value=screen["screen_label"])
        cell.font = Font(name="Arial", size=10, bold=True, color="333333")
        cell.alignment = Alignment(horizontal="left", vertical="center")
        current_row += 1

        # Header row
        ws.row_dimensions[current_row].height = 20
        for col, header in enumerate(["title", "audio", "ost(on screen text)", "graphics"], 1):
            cell = ws.cell(row=current_row, column=col, value=header)
            cell.fill = fill(HEADER_BG)
            cell.font = Font(name="Arial", size=10, bold=True, color=HEADER_FG)
            cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)
            cell.border = thin_border()
        current_row += 1

        # Data row
        ws.row_dimensions[current_row].height = 100
        for col, value in enumerate([screen["title"], audio, ost, ""], 1):
            cell = ws.cell(row=current_row, column=col, value=value)
            cell.font = Font(name="Arial", size=10, bold=(col == 1), color="1A1A2E")
            cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True, indent=1)
            cell.border = thin_border()
            if col == 4:
                cell.font = Font(name="Arial", size=10, color="999999", italic=True)
        current_row += 1

        # Blank row
        ws.row_dimensions[current_row].height = 10
        current_row += 1

    wb.save(output_path)
    print(f"✅ Saved: {output_path}")

if __name__ == "__main__":
    build_excel("storyboard_pro.pptx", "storyboard_script.xlsx")
